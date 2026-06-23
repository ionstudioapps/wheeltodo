"""WheelTodo — beta-launch positioning & GTM deck. Grounded in competitive-brief.md.
Built with python-pptx (no Node/LibreOffice available). Conservative layout, content-QA only.
"""
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# palette
INK   = RGBColor(0x14,0x12,0x10)
CORAL = RGBColor(0xFF,0x5C,0x4D)
CORALd= RGBColor(0xD8,0x42,0x35)
BEIGE = RGBColor(0xF4,0xF1,0xEC)
BEIGE2= RGBColor(0xEC,0xE6,0xDB)
WHITE = RGBColor(0xFF,0xFF,0xFF)
MUTE  = RGBColor(0x7A,0x70,0x66)
MUTEd = RGBColor(0xB9,0xAF,0xA3)
INKsoft = RGBColor(0x3A,0x35,0x30)
PURPLE= RGBColor(0xA7,0x8B,0xFA)
TEAL  = RGBColor(0x2B,0xB6,0xA3)
ICE   = RGBColor(0xEC,0xE9,0xE3)
F = "Arial"

prs = Presentation()
prs.slide_width  = I(13.333)
prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def _noshadow(sh):
    try: sh.shadow.inherit = False
    except Exception: pass

def rect(s,l,t,w,h,fill=None,line=None,lw=1.0,rounded=True,radius=0.08):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             I(l),I(t),I(w),I(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    _noshadow(shp)
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp

def oval(s,l,t,d,fill,line=None,lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL,I(l),I(t),I(d),I(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    _noshadow(shp); return shp

def tb(s,l,t,w,h,anchor=MSO_ANCHOR.TOP,wrap=True):
    box = s.shapes.add_textbox(I(l),I(t),I(w),I(h))
    tf = box.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf

def para(tf, runs, align=PP_ALIGN.LEFT, before=0, after=4, line=1.0, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if before: p.space_before = Pt(before)
    p.space_after = Pt(after)
    try: p.line_spacing = line
    except Exception: pass
    if isinstance(runs,str): runs = [(runs,16,INK,False,False,F)]
    for spec in runs:
        txt,size,color,bold = spec[0],spec[1],spec[2],spec[3]
        italic = spec[4] if len(spec)>4 else False
        fnt = spec[5] if len(spec)>5 else F
        r = p.add_run(); r.text = txt
        r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
        r.font.name=fnt; r.font.color.rgb=color
    return p

def kicker(s, text, color=CORAL, l=0.62, t=0.52):
    tf = tb(s,l,t,9,0.34)
    para(tf,[(text.upper(),12.5,color,True)],first=True,after=0)

def title(s, text, l=0.62, t=0.86, w=12.1, size=31, color=INK):
    tf = tb(s,l,t,w,1.0)
    para(tf,[(text,size,color,True)],first=True,after=0,line=1.02)

def cite(s, text, l=0.62, t=7.02, w=12.1, color=MUTE):
    tf = tb(s,l,t,w,0.34)
    para(tf,[(text,9,color,False,True)],first=True,after=0)

def wheelmotif(s, cx, cy, d, ring=CORAL, dot=WHITE, segs=BEIGE2):
    oval(s, cx-d/2, cy-d/2, d, ring)
    inner = d*0.40
    oval(s, cx-inner/2, cy-inner/2, inner, dot)

# ---------------------------------------------------------------- 1 · TITLE
s = slide(INK)
oval(s, 10.05, -1.7, 5.2, CORALd)         # offset glow
wheelmotif(s, 11.55, 2.05, 3.0)
tf = tb(s,0.9,2.35,8.6,1.4)
p = tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
for t_,c in [("Wheel",CORAL),("Todo",WHITE)]:
    r=p.add_run(); r.text=t_; r.font.size=Pt(58); r.font.bold=True; r.font.name=F; r.font.color.rgb=c
para(tf,[("Spin. Focus. Done.",22,ICE,False,True)],after=0,before=6)
tf2 = tb(s,0.92,4.35,8.7,1.5)
para(tf2,[("The kind to-do app that decides for you —",20,WHITE,True)],first=True,after=2,line=1.12)
para(tf2,[("and lets you rest without losing your streak.",20,CORAL,True)],after=0,line=1.12)
tf3 = tb(s,0.92,6.55,11,0.4)
para(tf3,[("Beta-launch positioning & go-to-market   ·   June 2026   ·   iOS · Android · Web",12,MUTEd,False)],first=True,after=0)

# ---------------------------------------------------------------- 2 · PROBLEM
s = slide(WHITE)
kicker(s,"The problem")
title(s,"You're not lazy. You're frozen.")
tf = tb(s,0.62,1.95,6.7,3.4)
para(tf,[("Decision paralysis is the real blocker — not laziness, not time.",17,INK,True)],first=True,after=10,line=1.15)
para(tf,[("A long list doesn't get done. Staring at 23 tasks, the overwhelmed (and especially ADHD) brain can't pick a starting point — so it picks nothing.",15,INKsoft,False)],after=10,line=1.25)
para(tf,[("The hardest part of any task is choosing it. So we delete the choosing.",15.5,CORALd,True,True)],after=0,line=1.2)
# stat cards
def stat(s,l,big,lab):
    rect(s,l,4.6,3.55,1.95,fill=BEIGE)
    tf=tb(s,l+0.3,4.82,3.0,0.95,anchor=MSO_ANCHOR.MIDDLE)
    para(tf,[(big,40,CORAL,True)],first=True,after=0)
    tf2=tb(s,l+0.3,5.78,3.0,0.62)
    para(tf2,[(lab,12.5,INKsoft,False)],first=True,after=0,line=1.05)
for l,big,lab in [(0.62,"6B+","views on #ADHD — TikTok"),
                  (4.5,"112M+","views on #productivitytips"),
                  (8.4,"~81%","completion on 0–10s clips")]:
    stat(s,l,big,lab)
# right rail
rect(s,7.55,1.95,5.18,2.35,fill=INK)
tf=tb(s,7.9,2.2,4.5,1.95,anchor=MSO_ANCHOR.MIDDLE)
para(tf,[("“Someone saying ‘this is my life’ —",15.5,WHITE,False,True)],first=True,after=2,line=1.2)
para(tf,[("not ‘here are the symptoms’ —",15.5,WHITE,False,True)],after=2,line=1.2)
para(tf,[("produces faster attitude change.”",15.5,CORAL,True,True)],after=6,line=1.2)
para(tf,[("First-person peer content out-performs explainers.",11.5,MUTEd,False)],after=0)
cite(s,"#ADHD / #productivitytips: Washington Post, 2022.  Clip-completion: SocialPilot, 2026.  Peer-vs-explainer: neurolaunch.com.")

# ---------------------------------------------------------------- 3 · AUDIENCE
s = slide(WHITE)
kicker(s,"Who it's for")
title(s,"Built for the brains that freeze")
cards = [("🧠","ADHD & neurodivergent","Task-initiation and decision paralysis are core. The true competitive neighborhood — and the most reachable audience on short-form."),
         ("🌊","The overwhelmed","Too much to do, so nothing gets done. Students, busy parents, anyone drowning in an endless list."),
         ("🎯","Focus-seekers","People who already love Pomodoro and ‘study-with-me’ — and want the deciding done for them.")]
cw, gap, l0, t0, ch = 3.86, 0.27, 0.62, 2.05, 3.55
for i,(em,h,body) in enumerate(cards):
    l = l0 + i*(cw+gap)
    rect(s,l,t0,cw,ch,fill=BEIGE)
    oval(s,l+0.32,t0+0.34,0.86,CORAL)
    tfi=tb(s,l+0.32,t0+0.34,0.86,0.86,anchor=MSO_ANCHOR.MIDDLE)
    para(tfi,[(em,26,WHITE,False)],first=True,after=0,align=PP_ALIGN.CENTER)
    tfh=tb(s,l+0.34,t0+1.45,cw-0.66,0.7)
    para(tfh,[(h,17,INK,True)],first=True,after=0,line=1.05)
    tfb=tb(s,l+0.34,t0+2.2,cw-0.66,1.2)
    para(tfb,[(body,13,INKsoft,False)],first=True,after=0,line=1.22)
tf=tb(s,0.62,5.95,12.1,0.9)
para(tf,[("Reachable wedge: ",14,INK,True),("lead with first-person ‘this fixed my task paralysis’ content, not feature explainers.",14,INKsoft,False)],first=True,after=0,line=1.2)
cite(s,"Neighborhood read & content guidance: competitive brief §1b, §3.")

# ---------------------------------------------------------------- 4 · WHAT IT IS
s = slide(WHITE)
kicker(s,"What it is")
title(s,"One tap from frozen to focused")
steps = [("1","Add your tasks","Dump everything onto the wheel. No sorting, no priority math."),
         ("2","Spin → ONE task","The wheel randomly lands on a single task. The choice is made for you."),
         ("3","Auto-Pomodoro","It drops you straight into a focus timer. You're already doing it.")]
cw,gap,l0,t0,ch = 3.86,0.27,0.62,2.0,2.15
for i,(n,h,body) in enumerate(steps):
    l=l0+i*(cw+gap)
    rect(s,l,t0,cw,ch,fill=INK)
    oval(s,l+0.3,t0+0.3,0.7,CORAL)
    tfn=tb(s,l+0.3,t0+0.3,0.7,0.7,anchor=MSO_ANCHOR.MIDDLE)
    para(tfn,[(n,22,WHITE,True)],first=True,after=0,align=PP_ALIGN.CENTER)
    tfh=tb(s,l+0.3,t0+1.12,cw-0.6,0.5)
    para(tfh,[(h,16.5,WHITE,True)],first=True,after=0)
    tfb=tb(s,l+0.3,t0+1.55,cw-0.6,0.5)
    para(tfb,[(body,12,MUTEd,False)],first=True,after=0,line=1.18)
# feature row
feats=[("🌿","Rest Mode","rest protects your streak"),("🔥","Streaks + 6 achievements","momentum that sticks"),
       ("🎨","4 accessible themes","light · dark · high-contrast"),("🔒","Offline-first","no account required")]
fw,fgap,fl0,ft0,fh=2.94,0.22,0.62,4.55,1.75
for i,(em,h,body) in enumerate(feats):
    l=fl0+i*(fw+fgap)
    rect(s,l,ft0,fw,fh,fill=BEIGE)
    tfi=tb(s,l+0.28,ft0+0.26,1.0,0.6)
    para(tfi,[(em,22,INK,False)],first=True,after=0)
    tfh=tb(s,l+0.28,ft0+0.92,fw-0.5,0.5)
    para(tfh,[(h,13.5,INK,True)],first=True,after=0,line=1.05)
    tfb=tb(s,l+0.28,ft0+1.36,fw-0.5,0.35)
    para(tfb,[(body,11,INKsoft,False)],first=True,after=0,line=1.05)
cite(s,"Feature set per docs/FEATURES.md (shipped build).  Body-double companion is roadmap — see final slide.")

# ---------------------------------------------------------------- 5 · WEDGE (hero, dark)
s = slide(INK)
oval(s,-1.6,5.2,4.6,CORALd)
wheelmotif(s,11.7,5.7,2.6,ring=CORAL)
kicker(s,"The wedge",color=CORAL,l=0.9,t=0.95)
tf=tb(s,0.9,1.7,11.4,2.4)
para(tf,[("Rest protects",62,WHITE,True)],first=True,after=0,line=1.0)
para(tf,[("your streak.",62,CORAL,True)],after=0,line=1.0)
tf2=tb(s,0.92,4.45,9.7,2.0)
para(tf2,[("Every other app treats rest as an ",18,ICE,False),("exemption",18,WHITE,True),
          (" — streak shields, vacation mode, ‘your streak pauses instead of resetting.’",18,ICE,False)],first=True,after=10,line=1.25)
para(tf2,[("WheelTodo is the only one where resting is a ",18,ICE,False),("positive move that keeps your streak alive.",18,CORAL,True),
          (" Rest Mode + the Drained / Okay / Restless check-in invert the category default.",18,ICE,False)],after=0,line=1.25)
cite(s,"Streak-as-damage-control framing: HabitStreak, Habi, Streaks (2026).  This position is currently unoccupied — competitive brief §2.",color=MUTEd)

# ---------------------------------------------------------------- 6 · COMPETITIVE MAP
s = slide(WHITE)
kicker(s,"Who we define against")
title(s,"Three foils, three weaknesses")
comp=[("Habitica","the punishment foil","Fail a daily and your avatar takes damage — critics say it ‘creates anxiety instead of motivation.’","“Your avatar doesn't die when you rest.”",CORAL),
      ("Finch","the gentle benchmark","~$30M ARR proving gentle works — but it's a self-care companion with no task-execution engine, no decision mechanic.","We add: actually start the task.",TEAL),
      ("Llama Life · Tiimo","the tonal cousins","Gentle and per-task-timed and genuinely good — but both still require you to decide what to start.","We remove the choice.",PURPLE)]
cw,gap,l0,t0,ch=3.86,0.27,0.62,2.0,4.35
for i,(name,tag,weak,counter,accent) in enumerate(comp):
    l=l0+i*(cw+gap)
    rect(s,l,t0,cw,ch,fill=BEIGE)
    oval(s,l+0.34,t0+0.36,0.34,accent)
    tfh=tb(s,l+0.84,t0+0.32,cw-1.1,0.5)
    para(tfh,[(name,17,INK,True)],first=True,after=0,line=0.98)
    tft=tb(s,l+0.34,t0+1.0,cw-0.66,0.4)
    para(tft,[(tag,12.5,accent if accent!=PURPLE else RGBColor(0x6B,0x4E,0xC2),True,True)],first=True,after=0)
    tfw=tb(s,l+0.34,t0+1.55,cw-0.66,1.95)
    para(tfw,[(weak,13,INKsoft,False)],first=True,after=0,line=1.24)
    rect(s,l+0.34,t0+3.45,cw-0.68,0.72,fill=INK)
    tfc=tb(s,l+0.5,t0+3.52,cw-1.0,0.6,anchor=MSO_ANCHOR.MIDDLE)
    para(tfc,[(counter,12.5,WHITE,True,True)],first=True,after=0,line=1.05)
cite(s,"Habitica critique: calmevo.com, 2026.  Finch ARR: blog.sparrowapps.io, 2025.  Full landscape: competitive brief §1.")

# ---------------------------------------------------------------- 7 · WHITE SPACE + honesty
s = slide(WHITE)
kicker(s,"Where we win")
title(s,"Two open lanes — and one honest caveat")
# two lanes
rect(s,0.62,2.0,5.85,2.5,fill=BEIGE)
tf=tb(s,0.92,2.24,5.3,2.1)
para(tf,[("1 · ‘Rest protects your streak’",16.5,INK,True)],first=True,after=4)
para(tf,[("Unoccupied. Streak-preservation always ships as damage control — never as a positive action you take. Ours is the inversion.",13,INKsoft,False)],after=0,line=1.24)
rect(s,6.68,2.0,6.05,2.5,fill=BEIGE)
tf=tb(s,6.98,2.24,5.5,2.1)
para(tf,[("2 · ‘Decision paralysis’ as a brand",16.5,INK,True)],first=True,after=4)
para(tf,[("The phrase lives in competitors' copy but almost no app is named or branded around it. High intent, low competition — ownable.",13,INKsoft,False)],after=0,line=1.24)
# honesty box
rect(s,0.62,4.75,12.11,1.95,fill=INK)
oval(s,1.0,5.05,0.62,CORAL)
tfi=tb(s,1.0,5.05,0.62,0.62,anchor=MSO_ANCHOR.MIDDLE)
para(tfi,[("!",26,WHITE,True)],first=True,after=0,align=PP_ALIGN.CENTER)
tf=tb(s,1.95,4.98,10.4,1.55,anchor=MSO_ANCHOR.MIDDLE)
para(tf,[("The wheel is NOT novel — don't claim invention.",15.5,CORAL,True)],first=True,after=3,line=1.1)
para(tf,[("Decision wheels are a saturated novelty category; two apps are literally named ‘Task Roulette.’ We own the ",13,ICE,False),
         ("integration",13,WHITE,True),
         (" — wheel → one task → auto-Pomodoro → Rest Mode — and the gentleness. Never the spinner itself.",13,ICE,False)],after=0,line=1.24)
cite(s,"‘Task Roulette’ apps: Google Play / monday.com, 2024.  Under-branding of ‘decision paralysis’: morgen.so, 2026.")

# ---------------------------------------------------------------- 8 · POSITIONING STATEMENT
s = slide(BEIGE)
kicker(s,"Positioning statement")
title(s,"How we say it, in one breath")
rect(s,0.62,2.05,12.11,3.25,fill=WHITE)
tf=tb(s,1.2,2.45,10.9,2.6,anchor=MSO_ANCHOR.MIDDLE)
para(tf,[("For ",20,INKsoft,False),("overwhelmed and ADHD brains who freeze at their to-do list",20,INK,True),
         (", WheelTodo is the ",20,INKsoft,False),("kind productivity app",20,CORALd,True),
         (" that ",20,INKsoft,False),("spins to pick one task and starts the timer for you",20,INK,True),
         (" — and ",20,INKsoft,False),("counts your rest",20,CORALd,True),(".",20,INKsoft,False)],first=True,after=12,line=1.3)
para(tf,[("Unlike Habitica, which punishes you for slipping — or planners that still make you choose.",16,INKsoft,False,True)],after=0,line=1.28)
tf2=tb(s,0.62,5.55,12.1,1.2)
for lab,val in [("Category line","‘the to-do app that decides for you’"),("Heart line","‘rest counts too’"),("Proof line","‘from 23-tasks-frozen to actually done’")]:
    pass
para(tf2,[("Category   ",12.5,CORALd,True),("the to-do app that decides for you        ",13,INK,False),
          ("Heart   ",12.5,CORALd,True),("rest counts too        ",13,INK,False),
          ("Proof   ",12.5,CORALd,True),("from 23-tasks-frozen to actually done",13,INK,False)],first=True,after=0,line=1.2)

# ---------------------------------------------------------------- 9 · HOOKS
s = slide(WHITE)
kicker(s,"Messaging · the hooks that sell it")
title(s,"Three hooks, ready to shoot")
hooks=[("Rest counts too","Collapse onto the couch, wrecked — phone shows the streak flame staying lit.","“I took a nap and my to-do app… congratulated me.”","the productivity app that rewards you for resting"),
       ("Spin → already doing it","Wheel mid-spin lands on one task → screen instantly flips to a running timer.","“Stop choosing. Spin — the timer's already counting.”","23 tasks. one spin. now you're already doing it."),
       ("It checks on you first","App asks ‘How's your energy?’ → tap Drained → it suggests gentle wins.","“It asks how I feel before it asks what I'll do.”","a to-do app that checks on you first")]
t0=2.0; rh=1.5; l0=0.62
# header row
hdr=["Hook","Visual (0–1s)","Verbal","Text overlay"]
xs=[0.62,3.0,6.7,9.9]; ws=[2.3,3.6,3.1,2.83]
tfh=tb(s,0,t0-0.42,13.333,0.36)
for x,w_,htxt in zip(xs,ws,hdr):
    b=s.shapes.add_textbox(I(x),I(t0-0.46),I(w_),I(0.36)).text_frame
    b.word_wrap=True; b.margin_left=0;b.margin_top=0;b.margin_bottom=0;b.margin_right=0
    para(b,[(htxt.upper(),11,CORALd,True)],first=True,after=0)
for i,(name,vis,verb,ov) in enumerate(hooks):
    t=t0+i*(rh+0.16)
    rect(s,0.62,t,12.11,rh,fill=BEIGE if i%2==0 else BEIGE)
    cols=[(0.86,2.1,[(name,15,INK,True)]),
          (3.0,3.5,[(vis,12.5,INKsoft,False)]),
          (6.7,3.05,[(verb,12.5,CORALd,True,True)]),
          (9.9,2.75,[(ov,12.5,INK,False)])]
    for x,w_,runs in cols:
        b=tb(s,x,t+0.16,w_,rh-0.3,anchor=MSO_ANCHOR.MIDDLE)
        para(b,runs,first=True,after=0,line=1.16)
cite(s,"Hooks map to the 30-day calendar (Hook Bank tab).  Caption every reel — most watch muted.")

# ---------------------------------------------------------------- 10 · GTM ENGINE
s = slide(WHITE)
kicker(s,"Go-to-market · the short-form engine")
title(s,"Short-form first, paid-aware")
# left: pillars
rect(s,0.62,2.0,6.0,4.55,fill=BEIGE)
tf=tb(s,0.92,2.24,5.5,4.1)
para(tf,[("6 content pillars (over 30 days)",15.5,INK,True)],first=True,after=8)
for nm,ct in [("Relatable pain","3"),("Differentiator: kind + rest","4"),("Oddly satisfying","4"),("Niche use-cases","6"),("Trend-jack","3"),("Proof & momentum","6")]:
    para(tf,[("•  ",13,CORAL,True),(nm,13,INK,False),("   —  "+ct+" posts",12,MUTE,False)],after=4,line=1.15)
para(tf,[("+ 4 on-brand rest days (story only). Full plan in the 30-day calendar.",11.5,MUTE,False,True)],before=4,after=0,line=1.18)
# right: reel formats
rect(s,6.78,2.0,5.95,4.55,fill=INK)
tf=tb(s,7.08,2.24,5.4,4.1)
para(tf,[("3 reel formats with proven fit",15.5,WHITE,True)],first=True,after=8)
for n,h,b in [("1","Satisfying spin-loop demo","spin → one task → timer-start. Screen-native, loopable, <10s."),
              ("2","First-person ‘this fixed my task paralysis’","peer POV beats explainers for this audience."),
              ("3","‘Apps that fixed my ADHD’ listicle","high reach via creator roundups & before/after.")]:
    para(tf,[(n+"  ",13,CORAL,True),(h,13.5,WHITE,True)],after=2,line=1.1)
    para(tf,[("     "+b,12,MUTEd,False)],after=7,line=1.16)
cite(s,"Formats & peer-content evidence: competitive brief §3 (neurolaunch.com; SocialPilot 2026; Structured/Finch case reads).")

# ---------------------------------------------------------------- 11 · ASO
s = slide(WHITE)
kicker(s,"App Store optimization")
title(s,"Own one term. Fight on three. Seed the rest.")
ladders=[("OWN it",CORAL,'“decision paralysis app”',"High intent, on-mechanic, rarely branded around. Our strongest niche term — anchor the subtitle here."),
         ("FIGHT on",INK,'“ADHD to-do” · “ADHD planner” · “ADHD focus timer”',"The battleground: high volume, high competition. Use for paid + long-tail (‘Pomodoro for ADHD’)."),
         ("SEED",TEAL,'“gentle productivity” · “task roulette” · “Pomodoro wheel”',"Low store volume, high SEO/PR & differentiation value. Nearly unclaimed — grab for content + defense.")]
t0=2.0; rh=1.42
for i,(tag,accent,terms,note) in enumerate(ladders):
    t=t0+i*(rh+0.18)
    rect(s,0.62,t,12.11,rh,fill=BEIGE)
    rect(s,0.62,t,2.3,rh,fill=accent)
    tft=tb(s,0.62,t,2.3,rh,anchor=MSO_ANCHOR.MIDDLE)
    para(tft,[(tag,16,WHITE,True)],first=True,after=0,align=PP_ALIGN.CENTER)
    tf=tb(s,3.1,t+0.2,9.4,rh-0.34,anchor=MSO_ANCHOR.MIDDLE)
    para(tf,[(terms,15,INK,True)],first=True,after=3,line=1.05)
    para(tf,[(note,12.5,INKsoft,False)],after=0,line=1.16)
cite(s,"Volume/competition reads are inferred from market structure — exact figures are paywalled in ASO tools.  competitive brief §4.")

# ---------------------------------------------------------------- 12 · CHANNELS + RISK
s = slide(WHITE)
kicker(s,"The honest truth & the big risk")
title(s,"Budget for paid. Watch retention, not spins.")
# left paid truth
rect(s,0.62,2.0,6.0,4.55,fill=INK)
tf=tb(s,0.95,2.3,5.4,4.0)
para(tf,[("The paid truth",16,CORAL,True)],first=True,after=8)
para(tf,[("Finch hit ~$30M ARR with zero VC — and even that ‘organic-looking’ growth is ",14,WHITE,False),("paid-driven.",14,CORAL,True)],after=8,line=1.26)
para(tf,[("Plan: budget paid TikTok / Meta from day one. Don't assume free virality. Community (Reddit / Discord) + paid amplification is the proven combo.",13.5,ICE,False)],after=0,line=1.26)
# right risk
rect(s,6.78,2.0,5.95,4.55,fill=BEIGE)
tf=tb(s,7.08,2.3,5.4,4.0)
para(tf,[("⚠  Riskiest assumption",16,CORALd,True)],first=True,after=8)
para(tf,[("That the wheel is a durable retention mechanic and not a one-week novelty.",14,INK,True)],after=8,line=1.24)
para(tf,[("Mitigation:",13.5,CORALd,True)],after=3)
for b in ["Make Rest Mode + streak — not the wheel — the retention spine.",
          "Instrument week-1 → week-4 spin frequency and rest-day usage.",
          "That signal is the primary beta success metric."]:
    para(tf,[("•  ",13,CORAL,True),(b,13,INKsoft,False)],after=4,line=1.18)
cite(s,"Finch growth model: blog.sparrowapps.io, 2025.  Risk analysis: competitive brief §5.")

# ---------------------------------------------------------------- 13 · ROADMAP + ASK (dark close)
s = slide(INK)
oval(s,10.2,4.9,5.4,CORALd)
wheelmotif(s,11.85,1.7,2.4)
kicker(s,"Roadmap & the ask",color=CORAL,l=0.9,t=0.7)
tf=tb(s,0.9,1.4,8.6,1.5)
para(tf,[("Ship the beta. Run the 30 days.",36,WHITE,True)],first=True,after=2,line=1.05)
para(tf,[("Watch the rest-streak signal.",36,CORAL,True)],after=0,line=1.05)
# roadmap card
rect(s,0.9,3.25,7.0,2.0,fill=INKsoft)
tf=tb(s,1.2,3.5,6.5,1.6)
para(tf,[("Next: the body-double companion  🌧️",15.5,CORAL,True)],first=True,after=5)
para(tf,[("A cozy lofi ‘study-with-me’ character (time-of-day × weather) living inside your tasks. Validated by gogh & Spirit City — but ",12.5,ICE,False),("roadmap, not shipped.",12.5,WHITE,True),(" Tease only.",12.5,ICE,False)],after=0,line=1.22)
# asks
tf=tb(s,8.25,3.25,4.4,3.2)
para(tf,[("The three deliverables",13.5,CORAL,True)],first=True,after=6)
for b in ["Competitive brief (cited)","30-day content calendar","This positioning deck"]:
    para(tf,[("✓  ",13,TEAL,True),(b,13.5,WHITE,False)],after=5,line=1.1)
tf=tb(s,0.9,5.65,11.5,1.2)
para(tf,[("Primary success metric:  ",14,CORAL,True),("week-1 → week-4 retention, led by rest-day usage — not spin count.",14,WHITE,True)],first=True,after=0,line=1.2)
cite(s,"Body-double benchmarks: gogh, Spirit City (Steam, 2026).  All claims sourced in 2026-06-23-competitive-brief.md.",color=MUTEd)

import os
out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"wheeltodo-launch-deck.pptx")
prs.save(out)
print("saved:",out,"| slides:",len(prs.slides._sldIdLst))
